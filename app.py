import os
import io
import json
import re
from flask import Flask, request, render_template, send_file
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)

# Use /tmp for file uploads on Vercel, as it's a writable directory
app.config['UPLOAD_FOLDER'] = '/tmp/uploads'
# Ensure the upload folder exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# --- Helper Functions ---

def text_to_slides_structure(text, guidance, api_key):
    """
    Uses the Gemini API to convert a block of text into a structured list of slides.
    """
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-1.5-flash') 
        
        prompt = f"""
        Analyze the following text and structure it into a series of presentation slides.
        Your response MUST be a valid JSON array of objects. Each object should represent a slide and have two keys: "title" (a string) and "content" (an array of strings, where each string is a bullet point).
        Do not include any explanatory text, markdown formatting, or code fences like ```json around the JSON output. The response should start with '[' and end with ']'.

        Guidance for tone/structure: "{guidance}"

        Text to process:
        ---
        {text}
        ---

        Example of desired JSON output:
        [
            {{"title": "Slide 1 Title", "content": ["Bullet point 1.", "Bullet point 2."]}},
            {{"title": "Slide 2 Title", "content": ["Another point.", "More details here."]}}
        ]
        """
        
        response = model.generate_content(prompt)
        
        json_text = re.search(r'\[.*\]', response.text, re.DOTALL)
        if not json_text:
            print("Error: No valid JSON array found in the Gemini response.")
            print("Raw Response:", response.text)
            return None

        structured_slides = json.loads(json_text.group(0))
        return structured_slides

    except Exception as e:
        print(f"An error occurred while calling the Gemini API: {e}")
        return None

def fallback_text_splitter(text):
    """A simple fallback to split text by paragraphs if the LLM fails."""
    paragraphs = text.strip().split('\n\n')
    slides = []
    for i, para in enumerate(paragraphs):
        lines = para.split('\n')
        title = lines[0]
        content = lines[1:] if len(lines) > 1 else [""]
        slides.append({"title": f"Slide {i+1}: {title}", "content": content})
    return slides


def create_presentation(slides_data, template_path):
    """
    Creates a PowerPoint presentation from structured slide data using a template.
    """
    prs = Presentation(template_path)
    
    title_slide_layout = prs.slide_layouts[0] 
    content_slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    for i, slide_info in enumerate(slides_data):
        layout = title_slide_layout if i == 0 else content_slide_layout
        slide = prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_info.get('title', 'Untitled Slide')

        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                body_shape = shape
                break
        
        if not body_shape:
            for shape in slide.placeholders:
                if ('content' in shape.name.lower() or 'body' in shape.name.lower()) and shape.placeholder_format.idx != 0:
                    body_shape = shape
                    break
        
        if not body_shape:
            for shape in slide.placeholders:
                if shape != title_shape and shape.has_text_frame:
                    body_shape = shape
                    break
        
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            content_list = slide_info.get('content', [])
            
            if isinstance(content_list, list):
                for point in content_list:
                    p = tf.add_paragraph()
                    p.text = str(point)
                    p.level = 0
            elif content_list:
                 p = tf.add_paragraph()
                 p.text = str(content_list)
                 p.level = 0
        else:
             print(f"Warning: Could not find a content placeholder for slide {i+1}.")

    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream

# --- Flask Routes ---

@app.route('/')
def index():
    """Renders the main page."""
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_presentation():
    """Handles the form submission and generates the presentation."""
    text = request.form.get('text')
    guidance = request.form.get('guidance')
    template_file = request.files.get('template')
    
    api_key = os.getenv("GEMINI_API_KEY")

    if not all([text, template_file, api_key]):
        missing = [name for name, var in {'text': text, 'template': template_file, 'API key': api_key}.items() if not var]
        return f"Missing required data: {', '.join(missing)}.", 400

    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_file.filename)
    template_file.save(template_path)

    slides_structure = text_to_slides_structure(text, guidance, api_key)
    
    if not slides_structure:
        slides_structure = fallback_text_splitter(text)
        if not slides_structure:
             os.remove(template_path)
             return "Could not generate slide structure from the text.", 500

    presentation_stream = create_presentation(slides_structure, template_path)
    os.remove(template_path)

    return send_file(
        presentation_stream,
        as_attachment=True,
        download_name='generated_presentation.pptx',
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

# The following check is not needed for Vercel, but good for local testing
if __name__ == '__main__':
    app.run(debug=True)

